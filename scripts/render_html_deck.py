#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""HTML 슬라이드 원본(play.html) → PPTX 재렌더링 파이프라인.

play.html 안의 `SLIDES` 배열(각 원소가 완결형 HTML 문서, assets/ 상대참조)을
헤드리스 Chrome으로 1280x720에서 렌더링하여 슬라이드별 PNG로 캡처하고,
16:9(13.333x7.5in) PPTX에 전면(full-bleed) 이미지로 배치한다.

사용 예:
    python scripts/render_html_deck.py                 # 기본 경로로 재렌더링
    python scripts/render_html_deck.py --src <play.html> --out <out.pptx> --scale 2

의존성: playwright(+시스템 Chrome), python-pptx, Pillow
    python -m pip install playwright python-pptx pillow
    # 브라우저는 시스템 Chrome/Edge를 사용하므로 `playwright install` 불필요
"""
import argparse
import json
import re
import sys
import tempfile
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

# --- 기본 경로 (이 환경 기준; 필요 시 --src/--out 로 덮어씀) ---
DEFAULT_SRC = Path(
    r"C:\Users\dnslab\Downloads\WageGuard 이주노동자 임금 착취 스크리닝\play.html"
)
DEFAULT_OUT = DEFAULT_SRC.parent / "WageGuard_발표_재렌더.pptx"

SLIDE_W_IN = 13.333  # 16:9
SLIDE_H_IN = 7.5
VIEW_W = 1280
VIEW_H = 720


def extract_slides(src: Path):
    """play.html에서 SLIDES 배열을 파싱해 HTML 문자열 리스트로 반환."""
    t = src.read_text(encoding="utf-8")
    m = re.search(r"(?:var|const|let)\s+SLIDES\s*=\s*(\[.*?\]);", t, re.S)
    if not m:
        raise SystemExit("SLIDES 배열을 찾지 못했습니다: " + str(src))
    slides = json.loads(m.group(1))
    if not slides:
        raise SystemExit("SLIDES 배열이 비어 있습니다.")
    return slides


def render_pngs(src: Path, slides, out_dir: Path, scale: int, settle_ms: int, channel: str):
    """각 슬라이드를 assets/ 옆에 임시 저장해 렌더링, PNG 경로 리스트 반환."""
    src_dir = src.parent
    png_paths = []
    with sync_playwright() as p:
        browser = p.chromium.launch(channel=channel, headless=True)
        page = browser.new_page(
            viewport={"width": VIEW_W, "height": VIEW_H},
            device_scale_factor=scale,
        )
        for i, html in enumerate(slides, 1):
            # assets/ 상대참조가 풀리도록 소스 폴더 안에 임시 파일로 저장
            tmp = src_dir / f"_render_tmp_{i:02d}.html"
            tmp.write_text(html, encoding="utf-8")
            try:
                page.goto(tmp.as_uri(), wait_until="networkidle")
                # 폰트 로딩 완료 대기
                try:
                    page.evaluate("document.fonts && document.fonts.ready")
                except Exception:
                    pass
                # gsap 애니메이션이 있으면 즉시 최종 상태로 완료
                page.evaluate(
                    "() => { try { if (window.gsap && gsap.globalTimeline) "
                    "gsap.globalTimeline.progress(1); } catch(e){} }"
                )
                page.wait_for_timeout(settle_ms)  # 잔여 전환/폰트 리플로우 여유
                png = out_dir / f"slide{i:02d}.png"
                page.screenshot(path=str(png), clip={"x": 0, "y": 0, "width": VIEW_W, "height": VIEW_H})
                png_paths.append(png)
                print(f"  렌더 완료 S{i:02d} -> {png.name}")
            finally:
                tmp.unlink(missing_ok=True)
        browser.close()
    return png_paths


def build_pptx(png_paths, out: Path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]  # 빈 레이아웃
    for png in png_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))


def main():
    ap = argparse.ArgumentParser(description="play.html → PPTX 재렌더링")
    ap.add_argument("--src", type=Path, default=DEFAULT_SRC, help="play.html 경로")
    ap.add_argument("--out", type=Path, default=DEFAULT_OUT, help="출력 pptx 경로")
    ap.add_argument("--scale", type=int, default=2, help="해상도 배수 (2 → 2560x1440)")
    ap.add_argument("--settle-ms", type=int, default=600, help="렌더 후 대기(ms)")
    ap.add_argument("--channel", default="chrome", help="브라우저 채널 (chrome/msedge)")
    ap.add_argument("--keep-png", action="store_true", help="PNG 중간산출물 보존")
    args = ap.parse_args()

    if not args.src.exists():
        raise SystemExit(f"소스 없음: {args.src}")

    print(f"[1/3] 슬라이드 추출: {args.src}")
    slides = extract_slides(args.src)
    print(f"      슬라이드 {len(slides)}장")

    png_dir = Path(tempfile.mkdtemp(prefix="wg_render_")) if not args.keep_png \
        else (args.out.parent / "_render_png")
    png_dir.mkdir(parents=True, exist_ok=True)
    print(f"[2/3] 렌더링 (Chrome headless, x{args.scale}) → {png_dir}")
    pngs = render_pngs(args.src, slides, png_dir, args.scale, args.settle_ms, args.channel)

    print(f"[3/3] PPTX 조립 → {args.out}")
    build_pptx(pngs, args.out)
    print(f"완료: {args.out}  ({len(pngs)}장)")
    if args.keep_png:
        print(f"PNG 보존: {png_dir}")


if __name__ == "__main__":
    main()
