# -*- coding: utf-8 -*-
"""
WageGuard 최종 발표(10분) PPTX 생성기.

python scripts/build_presentation.py
→ WageGuard_발표.pptx (16:9)

발표자: 선문대학교 컴퓨터공학과 김태성
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# 디자인 시스템
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x0F, 0x2A, 0x43)      # 딥 네이비 (표지/헤더)
NAVY2 = RGBColor(0x1B, 0x3A, 0x5A)
TEAL = RGBColor(0x14, 0xB8, 0xA6)      # 강조 (틸)
TEAL_D = RGBColor(0x0D, 0x9488 // 256, 0x9488 % 256)  # 사용 안함 방지용
AMBER = RGBColor(0xF5, 0x9E, 0x0B)     # 보조 강조 (앰버)
CORAL = RGBColor(0xEF, 0x4B, 0x4B)     # 위험/문제
INK = RGBColor(0x1E, 0x29, 0x3B)       # 본문 진한 글자
GRAY = RGBColor(0x64, 0x74, 0x8B)      # 회색 본문
LGRAY = RGBColor(0x94, 0xA3, 0xB8)
BG = RGBColor(0xF7, 0xF9, 0xFC)        # 밝은 배경
CARD = RGBColor(0xFF, 0xFF, 0xFF)
CARD_LINE = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHIP_BG = RGBColor(0xE6, 0xF7, 0xF4)

FONT = "맑은 고딕"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG_DIR = os.path.join(HERE, "results", "_backup_v1")


def img(name: str) -> str | None:
    for d in (os.path.join(HERE, "results"), IMG_DIR):
        p = os.path.join(d, name)
        if os.path.exists(p):
            return p
    return None


# ----------------------------------------------------------------------------
# 저수준 헬퍼
# ----------------------------------------------------------------------------
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=None, line_w=None, shadow=False,
         rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=Pt(4), line_spacing=1.05, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = FONT
    return tb


def bg(slide, color=BG):
    rect(slide, 0, 0, EMU_W, EMU_H, color)


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def header(slide, kicker, title, page, total=17):
    """콘텐츠 슬라이드 상단 헤더."""
    bg(slide, BG)
    # 상단 강조 바
    rect(slide, 0, 0, Inches(0.22), EMU_H, TEAL)
    # kicker (섹션 라벨)
    txt(slide, Inches(0.75), Inches(0.45), Inches(11), Inches(0.35),
        [[(kicker, 13, TEAL, True)]])
    # 제목
    txt(slide, Inches(0.72), Inches(0.75), Inches(11.8), Inches(0.9),
        [[(title, 30, NAVY, True)]])
    # 구분선
    rect(slide, Inches(0.75), Inches(1.62), Inches(11.85), Pt(2.2), NAVY)
    # 페이지
    txt(slide, Inches(11.9), Inches(6.95), Inches(1.2), Inches(0.35),
        [[(f"{page:02d} / {total:02d}", 10, LGRAY, False)]], align=PP_ALIGN.RIGHT)
    # 워터마크 풋터
    txt(slide, Inches(0.75), Inches(6.95), Inches(6), Inches(0.35),
        [[("WageGuard · 셀 단위 감독 우선순위 스크리닝", 10, LGRAY, False)]])


def chip(slide, x, y, w, text, color=TEAL, fill=CHIP_BG):
    c = rect(slide, x, y, w, Inches(0.42), fill, rounded=True)
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = FONT
    return c


def stat_card(slide, x, y, w, h, big, label, sub=None, accent=TEAL):
    rect(slide, x, y, w, h, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
    rect(slide, x, y, Inches(0.12), h, accent, rounded=False)
    runs = [[(big, 30, accent, True)], [(label, 12.5, INK, True)]]
    if sub:
        runs.append([(sub, 10.5, GRAY, False)])
    txt(slide, x + Inches(0.28), y + Inches(0.22), w - Inches(0.4), h - Inches(0.3),
        runs, anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(3), line_spacing=1.0)


# ----------------------------------------------------------------------------
# 슬라이드들
# ----------------------------------------------------------------------------
def s01_title(prs):
    s = add_slide(prs)
    bg(s, NAVY)
    # 배경 장식
    rect(s, 0, Inches(6.9), EMU_W, Inches(0.6), NAVY2)
    rect(s, Inches(0.9), Inches(1.55), Inches(1.5), Pt(4), TEAL)
    # 카테고리
    txt(s, Inches(0.95), Inches(1.0), Inches(11), Inches(0.4),
        [[("제5회 고용노동 공공데이터·AI 활용 공모전  |  아이디어 기획 부문", 14, TEAL, True)]])
    # 타이틀
    txt(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(1.4),
        [[("WageGuard", 60, WHITE, True)]])
    txt(s, Inches(0.95), Inches(3.15), Inches(11.6), Inches(1.4),
        [[("공공데이터·AI 기반 이주노동자", 30, WHITE, True)],
         [("임금 착취 ", 30, WHITE, True), ("위험 스크리닝", 30, TEAL, True),
          (" · 감독 우선순위 제안 시스템", 30, WHITE, True)]],
        line_spacing=1.15)
    # 부제
    txt(s, Inches(0.95), Inches(5.0), Inches(11.6), Inches(0.6),
        [[("신고 이전 단계에서, 신고할 수 없는 피해자를 국가가 먼저 찾아가도록", 15, LGRAY, False)]])
    # 발표자
    rect(s, Inches(0.95), Inches(5.95), Inches(5.2), Pt(1.5), NAVY2)
    txt(s, Inches(0.95), Inches(6.1), Inches(11), Inches(0.6),
        [[("선문대학교 컴퓨터공학과  ", 15, WHITE, True), ("김태성", 15, TEAL, True),
          ("      github.com/kts6450/WageGuard", 12, LGRAY, False)]])


def s02_agenda(prs):
    s = add_slide(prs)
    header(s, "AGENDA", "발표 순서", 2)
    items = [
        ("01", "문제 정의", "신고할 수 없는 111만 외국인 취업자"),
        ("02", "해결 아이디어", "WageGuard — 셀 단위 스크리닝 파이프라인"),
        ("03", "데이터 & AI", "공공데이터 3계층 + 4종 AI 모델"),
        ("04", "실증 & 검증", "공개 데이터 재구축 · 외부 명단 5.2배 포착"),
        ("05", "기대효과", "점검 효율 개선 · 찾아가는 행정"),
        ("06", "한계 · 로드맵", "정직한 한계 인정 · 단계적 실증"),
    ]
    x0, y0 = Inches(0.9), Inches(2.15)
    cw, ch, gx, gy = Inches(5.85), Inches(1.28), Inches(0.35), Inches(0.32)
    for i, (n, t, d) in enumerate(items):
        col, row = i % 2, i // 2
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        rect(s, x, y, cw, ch, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
        rect(s, x, y, Inches(0.14), ch, TEAL, rounded=False)
        txt(s, x + Inches(0.35), y + Inches(0.16), Inches(1.2), Inches(1),
            [[(n, 30, TEAL, True)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(1.55), y + Inches(0.2), cw - Inches(1.7), Inches(1),
            [[(t, 18, NAVY, True)], [(d, 12, GRAY, False)]],
            anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(3))


def s03_problem(prs):
    s = add_slide(prs)
    header(s, "01  PROBLEM", "신고할 수 없는 피해자들", 3)
    # 리드 문장
    txt(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(0.7),
        [[("이주노동자 임금 착취 ", 18, INK, False), ("신고율은 5% 미만", 18, CORAL, True),
          (". 가장 취약한 사람이 가장 신고하지 못한다.", 18, INK, False)]])
    # 통계 카드 3개
    y = Inches(2.75)
    cw, ch, gx = Inches(3.7), Inches(1.75), Inches(0.28)
    x = Inches(0.9)
    stat_card(s, x, y, cw, ch, "110.9만 명", "국내 외국인 취업자 · 역대 최고",
              "광·제조업 46만 · E-9 32만 (통계청 2025)", accent=NAVY)
    stat_card(s, x + (cw + gx), y, cw, ch, "계약 ≠ 현실", "표준계약 체결에도 위반 대응 곤란",
              "언어 장벽·체류자격(비자) 종속", accent=AMBER)
    stat_card(s, x + 2 * (cw + gx), y, cw, ch, "< 5%", "실제 신고 도달률",
              "언어장벽·비자취소 위협·감독 한계", accent=CORAL)
    # 원인 박스
    y2 = Inches(4.85)
    rect(s, Inches(0.9), y2, Inches(11.6), Inches(1.75), RGBColor(0xFD, 0xF2, 0xF2),
         line=RGBColor(0xF7, 0xC9, 0xC9), line_w=Pt(1), rounded=True)
    txt(s, Inches(1.2), y2 + Inches(0.2), Inches(11), Inches(0.4),
        [[("신고에 도달하지 못하는 구조적 이유", 15, CORAL, True)]])
    reasons = [
        "언어 장벽으로 신고 절차 접근성 부족",
        "체류자격(비자) 취소 위협에 대한 두려움",
        "근로감독관 1인당 담당 사업장 수천 개 — 구조적 자원 한계",
    ]
    txt(s, Inches(1.2), y2 + Inches(0.68), Inches(11), Inches(1),
        [[("•  " + r, 13.5, INK, False)] for r in reasons],
        space_after=Pt(6), line_spacing=1.0)


def s04_gap(prs):
    s = add_slide(prs)
    header(s, "01  PROBLEM", "기존 접근의 한계", 4)
    txt(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(0.6),
        [[("지금의 대응은 ", 17, INK, False), ("피해자의 자력 신고", 17, CORAL, True),
          ("에 의존한다.", 17, INK, False)]])
    rows = [
        ("신고 기반 사후 대응", "피해자 자력 신고에 의존 → 가장 취약한 이주노동자가 가장 적게 신고"),
        ("단순 통계 보고서", "산업·연도 단위 거시 추이 중심 → 현장 점검 우선순위로 환원되지 않음"),
        ("기존 부정수급 탐지", "사업주 행정 데이터 위주 → 임금·근로조건 위반의 미시 패턴 미반영"),
    ]
    y = Inches(2.7)
    ch, gy = Inches(1.15), Inches(0.25)
    for i, (t, d) in enumerate(rows):
        yy = y + i * (ch + gy)
        rect(s, Inches(0.9), yy, Inches(11.6), ch, CARD, line=CARD_LINE,
             line_w=Pt(1), rounded=True)
        rect(s, Inches(0.9), yy, Inches(3.6), ch, NAVY, rounded=False)
        txt(s, Inches(1.15), yy, Inches(3.2), ch, [[(t, 16, WHITE, True)]],
            anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(4.75), yy, Inches(7.5), ch, [[(d, 13.5, INK, False)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)


def s05_idea(prs):
    s = add_slide(prs)
    header(s, "02  SOLUTION", "WageGuard — 어떻게 작동하는가", 5)
    txt(s, Inches(0.9), Inches(1.85), Inches(11.6), Inches(0.7),
        [[("공공데이터를 ", 17, INK, False), ("AI 위험 스코어", 17, TEAL, True),
          ("로 바꿔, 신고 이전에 ", 17, INK, False),
          ("산업×지역×고용형태 셀 단위 감독 우선순위", 17, TEAL, True),
          ("를 제안한다.", 17, INK, False)]])
    # 3단계 파이프라인
    steps = [
        ("입력", "전면 공개 행정데이터",
         "국민연금 사업장 54만\n임금체불 공개 명단\nE-9 근무현황 (승인 불필요)", NAVY),
        ("1  셀 프로파일링", "실측 위험 신호 3종",
         "저임금(인당 보험료)\n고용 불안정(취득·상실)\n영세성(5인 미만 비중)", TEAL),
        ("2  AI 스코어링", "4종 모델 앙상블",
         "XGBoost · LNN/LSTM\nPyOD 이상탐지\n+ 이주노동자 가중치", AMBER),
        ("3  대시보드", "셀 단위 우선순위 산출",
         "산업×지역×고용형태\n위험 지도 · 우선순위\nPrecision@K 외부 검증", NAVY),
    ]
    x0, y = Inches(0.9), Inches(2.85)
    cw, ch = Inches(2.72), Inches(2.5)
    gap = Inches(0.28)
    for i, (tag, title, body, col) in enumerate(steps):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
        rect(s, x, y, cw, Inches(0.55), col, rounded=False)
        txt(s, x + Inches(0.15), y, cw - Inches(0.3), Inches(0.55),
            [[(tag, 13, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.2), y + Inches(0.7), cw - Inches(0.4), Inches(0.6),
            [[(title, 14.5, NAVY, True)]], align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.2), y + Inches(1.35), cw - Inches(0.4), Inches(1.1),
            [[(ln, 11.5, GRAY, False)] for ln in body.split("\n")],
            align=PP_ALIGN.CENTER, space_after=Pt(2), line_spacing=1.0)
        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                     x + cw + Emu(int(gap) // 2) - Inches(0.16),
                                     y + ch / 2 - Inches(0.16),
                                     Inches(0.32), Inches(0.32))
            _set_fill(arr, TEAL)
    # 결과 배너
    yb = Inches(5.75)
    rect(s, Inches(0.9), yb, Inches(11.6), Inches(0.85), NAVY, rounded=True)
    txt(s, Inches(1.2), yb, Inches(11), Inches(0.85),
        [[("결과 → ", 15, TEAL, True),
          ("셀 단위 감독 우선순위 · 이주노동자 지원 선제 배정  ", 15, WHITE, True),
          ("(개별 사업장 지목·자동 처벌 아님, 감독관 판단 보조)", 12, LGRAY, False)]],
        anchor=MSO_ANCHOR.MIDDLE)


def s06_data(prs):
    s = add_slide(prs)
    header(s, "03  DATA", "활용 공공데이터 — 3계층", 6)
    cols = [
        ("고용노동부 (주최)", NAVY,
         ["임금체불 사업주 명단 → 외부 검증",
          "연도별 근로사건 현황",
          "고용형태별근로실태조사 개방판",
          "최저임금위원회 (라벨 근거)"]),
        ("한국고용정보원 (주관)", TEAL,
         ["외국인근로자 근무현황",
          "→ 이주노동자 가중치",
          "EIS 고용행정통계 (노출 보정)",
          "워크피디아 임금분포 (보조 라벨)"]),
        ("공개 행정데이터 (실증 실사용)", AMBER,
         ["국민연금 사업장 내역 54만 곳",
          "임금체불 공개 명단 769명(중복 제거)",
          "E-9 시도×업종 28.5만 (KOSIS)",
          "(보조 연구) MDIS 491만 건"]),
    ]
    x0, y = Inches(0.9), Inches(2.05)
    cw, ch, gap = Inches(3.75), Inches(3.15), Inches(0.22)
    for i, (title, col, items) in enumerate(cols):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
        rect(s, x, y, cw, Inches(0.7), col, rounded=False)
        txt(s, x + Inches(0.2), y, cw - Inches(0.4), Inches(0.7),
            [[(title, 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.28), y + Inches(0.9), cw - Inches(0.5), ch - Inches(1),
            [[("• " + it, 12.5, INK, False)] for it in items],
            space_after=Pt(9), line_spacing=1.05)
    # 가점 배너
    yb = Inches(5.55)
    rect(s, Inches(0.9), yb, Inches(11.6), Inches(1.05), CHIP_BG,
         line=TEAL, line_w=Pt(1.2), rounded=True)
    txt(s, Inches(1.2), yb + Inches(0.12), Inches(11), Inches(0.9),
        [[("가점 “나” 충족  ", 15, TEAL, True),
          ("한국고용정보원(주관) 공공데이터 직접 활용 → +2점", 14, INK, True)],
         [("개인 식별 변수 미사용 · MDIS 비밀유지 서약 준수 · 개방판과 변수 정합 유지", 12, GRAY, False)]],
        space_after=Pt(5))


def s07_label(prs):
    s = add_slide(prs)
    header(s, "03  DATA", "문제의 규모 — 법 기준 위반 노출률 실측 (MDIS 보조 연구)", 7)
    conds = [
        ("A", "4대보험 미가입", "고용·건강·국민연금 중\n하나라도 미가입", "5.4%", NAVY),
        ("B", "최저임금 미달", "시간당임금 < 연도별\n법정 최저임금", "1.4%", AMBER),
        ("C", "초과근무 착취", "초과근무 20h+ 이면서\n초과급여 = 0", "0.1%", CORAL),
    ]
    x0, y = Inches(0.9), Inches(2.1)
    cw, ch, gap = Inches(3.1), Inches(2.5), Inches(0.3)
    for i, (a, t, d, rate, col) in enumerate(conds):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.28),
                                  y + Inches(0.25), Inches(0.6), Inches(0.6))
        _set_fill(circ, col)
        tf = circ.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = a
        r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
        r.font.name = FONT
        txt(s, x + Inches(0.28), y + Inches(1.0), cw - Inches(0.5), Inches(0.5),
            [[(t, 16, NAVY, True)]])
        txt(s, x + Inches(0.28), y + Inches(1.5), cw - Inches(0.5), Inches(0.7),
            [[(ln, 12, GRAY, False)] for ln in d.split("\n")],
            space_after=Pt(1), line_spacing=1.0)
        txt(s, x + Inches(0.28), y + Inches(1.95), cw - Inches(0.5), Inches(0.45),
            [[(rate + " 해당", 15, col, True)]])
    # 합집합 결과
    yb = Inches(4.95)
    rect(s, Inches(0.9), yb, Inches(11.6), Inches(1.55), NAVY, rounded=True)
    txt(s, Inches(1.3), yb, Inches(4.5), Inches(1.55),
        [[("6.7%", 46, TEAL, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.3), yb, Inches(7.8), Inches(1.55),
        [[("법 기준 위반 노출률 (A ∪ B ∪ C 합집합, 실측)", 17, WHITE, True)],
         [("491만 6,667건 중 실측 · 규칙 기반 라벨의 순환논리·누수 방어는 11p에서 선제 검증", 13, LGRAY, False)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(6))


def s08_ai(prs):
    s = add_slide(prs)
    header(s, "04  AI", "아키텍처 — 검증이 먼저, 고도화는 그 위에", 7)
    cards = [
        ("① 신호 엔진", "가동 중 — 3신호 합성", "저임금·고이직·영세성을\nz-표준화 합성 → 셀 위험점수", NAVY),
        ("② 검증 엔진", "가동 중 — 외부 채점", "체불 명단 독립 채점 5.2배\n순열검정·가중치 민감도", TEAL),
        ("③ 학습 확장", "로드맵 — XGBoost", "점검 결과(명중률) 축적 시\n데이터가 가중치를 학습", AMBER),
        ("④ 시계열 확장", "로드맵 — LNN 경보", "고용보험 이력 연계 시\n일 단위 준실시간 경보", CORAL),
    ]
    x0, y = Inches(0.9), Inches(2.1)
    cw, ch, gap = Inches(2.72), Inches(2.35), Inches(0.28)
    for i, (t, sub, body, col) in enumerate(cards):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
        rect(s, x, y, cw, Inches(0.12), col, rounded=False)
        txt(s, x + Inches(0.22), y + Inches(0.3), cw - Inches(0.4), Inches(0.5),
            [[(t, 16, col, True)]])
        txt(s, x + Inches(0.22), y + Inches(0.85), cw - Inches(0.4), Inches(0.4),
            [[(sub, 13, NAVY, True)]])
        txt(s, x + Inches(0.22), y + Inches(1.35), cw - Inches(0.4), Inches(0.9),
            [[(ln, 11.5, GRAY, False)] for ln in body.split("\n")],
            space_after=Pt(2), line_spacing=1.0)
    # 하단 신뢰성 배너
    yb = Inches(4.85)
    rect(s, Inches(0.9), yb, Inches(11.6), Inches(1.65), CARD, line=CARD_LINE,
         line_w=Pt(1), rounded=True)
    txt(s, Inches(1.2), yb + Inches(0.18), Inches(11), Inches(0.45),
        [[("신뢰성 보강 — 심사·정책 활용을 위한 검증 구조", 15, TEAL, True)]])
    pts = [
        "누수 통제: 위험 지표(국민연금)와 채점 데이터(체불 명단)를 소스부터 분리 — 순환 불가 구조",
        "가중치: 임의 상수 대신 동일가중 z-합성 + 200조합 섭동 민감도(순위 상관 0.93) + 과거 차수 학습 계수",
        "외부 검증 완료: 임금체불 공개 명단 769명 매칭 → 상위 10% 셀 포착 30% (고용 기준 5.2배·사업장 기준 2.4배)",
    ]
    txt(s, Inches(1.2), yb + Inches(0.68), Inches(11.2), Inches(0.9),
        [[("•  " + p, 12.5, INK, False)] for p in pts],
        space_after=Pt(4), line_spacing=1.0)


def s09_product(prs):
    """아이디어 부문 = 시연 불가 → 스크린샷이 시연을 대신한다."""
    s = add_slide(prs)
    header(s, "04  PRODUCT", "아이디어가 아니라 실물 — 감독관용 운영 시스템", 8)
    im = img("ui_ops_app.png")
    if im:
        rect(s, Inches(0.9), Inches(1.95), Inches(7.6), Inches(4.55), CARD,
             line=CARD_LINE, line_w=Pt(1), rounded=True)
        s.shapes.add_picture(im, Inches(1.02), Inches(2.2), width=Inches(7.36))
    x = Inches(8.75)
    txt(s, x, Inches(2.0), Inches(3.75), Inches(0.6),
        [[("감독관 워크플로 (구현 완료)", 15, NAVY, True)]])
    flow = [
        ("① 관할 선택", "시군구×업종 1,873셀 중 내 관할"),
        ("② 우선순위 + 근거", "“고이직+저임금” 등 신호 명시"),
        ("③ 점검 계획서 발급", "CSV 한 클릭 — 현장 투입"),
        ("④ 결과 입력 → 환류", "명중률 자동 집계·지표 재보정"),
    ]
    y = Inches(2.6)
    for i, (t, d) in enumerate(flow):
        yy = y + i * Inches(0.82)
        rect(s, x, yy, Inches(3.75), Inches(0.7), CARD, line=CARD_LINE,
             line_w=Pt(1), rounded=True)
        rect(s, x, yy, Inches(0.12), Inches(0.7), TEAL)
        txt(s, x + Inches(0.28), yy + Inches(0.07), Inches(3.4), Inches(0.6),
            [[(t, 12.5, NAVY, True)], [(d, 10.5, GRAY, False)]],
            space_after=Pt(1), line_spacing=1.0)
    txt(s, x, Inches(6.0), Inches(3.75), Inches(0.6),
        [[("매월 자동 갱신 · 개별 사업장 자동 처분 없음", 10.5, GRAY, False)]],
        line_spacing=1.1)


def s10_results(prs):
    s = add_slide(prs)
    header(s, "05  RESULTS", "실증 결과 — 전면 공개 데이터로 재구축·검증 완료", 9)
    cards = [
        ("54만 곳", "국민연금 전 사업장 실측", "2026-05 · 가입자 1,118만 명", NAVY),
        ("769명", "임금체불 공개 명단(중복 제거)", "2023~2026 6개 차수 전수 수집", TEAL),
        ("5.2배", "상위 10% 셀 체불 포착 (고용 기준)", "고용 6% → 체불 30% · 사업장 수 기준 2.4배", AMBER),
        ("4.7배", "홀드아웃 차수(25~26) 검증", "지표 산출에 전혀 쓰지 않은 명단", CORAL),
    ]
    x0, y = Inches(0.9), Inches(2.0)
    cw, ch, gap = Inches(2.72), Inches(1.6), Inches(0.28)
    for i, (b, l, sub, col) in enumerate(cards):
        x = x0 + i * (cw + gap)
        stat_card(s, x, y, cw, ch, b, l, sub, accent=col)
    # 검증 구조 배너 — 순환논리가 구조적으로 불가능한 이유
    yb = Inches(3.95)
    rect(s, Inches(0.9), yb, Inches(11.6), Inches(1.15), NAVY, rounded=True)
    txt(s, Inches(1.2), yb + Inches(0.14), Inches(11), Inches(0.9),
        [[("순환논리가 구조적으로 불가능한 검증", 14.5, TEAL, True)],
         [("위험 지표는 국민연금 행정데이터(저임금·고이직·영세성)만으로 산출 — "
           "임금체불 명단은 오직 채점에만 사용. 서로 독립된 두 공공데이터.",
           12.5, RGBColor(0xE2, 0xE8, 0xF0), False)]],
        space_after=Pt(5), line_spacing=1.05)
    # 하단 2카드 — 가중치 민감도 / 시간 분할 검증
    y2 = Inches(5.35)
    cw2 = Inches(5.66)
    rect(s, Inches(0.9), y2, cw2, Inches(1.25), CARD, line=CARD_LINE,
         line_w=Pt(1), rounded=True)
    txt(s, Inches(1.15), y2 + Inches(0.15), cw2 - Inches(0.5), Inches(1),
        [[("가중치 민감도 검증", 13.5, TEAL, True)],
         [("동일가중 대신 무작위 가중 200조합으로 바꿔도 셀 순위 상관 평균 0.93 — "
           "임의 가중치 논란 차단", 12, INK, False)]],
        space_after=Pt(4), line_spacing=1.05)
    x2 = Inches(0.9) + cw2 + Inches(0.28)
    rect(s, x2, y2, cw2, Inches(1.25), CARD, line=CARD_LINE,
         line_w=Pt(1), rounded=True)
    txt(s, x2 + Inches(0.25), y2 + Inches(0.15), cw2 - Inches(0.5), Inches(1),
        [[("우연도, 업종 몰빵도 아님", 13.5, TEAL, True)],
         [("순열검정 p=0.006 (무작위 1만 회 대비) · 건설업을 통째로 빼도 "
           "4.0배(고용)/2.9배(사업장) 유지", 12, INK, False)]],
        space_after=Pt(4), line_spacing=1.05)


def s10b_defense(prs):
    s = add_slide(prs)
    header(s, "05  VALIDATION", "심사 전에, 스스로에게 먼저 던진 세 가지 질문", 10)
    rows = [
        ("순환논리 아닌가?",
         "규칙으로 라벨을 만들고 같은 규칙을 예측하면 성능이 과대평가된다",
         "재구축 검증은 지표(국민연금)와 채점(체불 명단)이 서로 독립된 데이터 — 구조적으로 순환 불가능. 홀드아웃 차수에서도 4.7배 유지"),
        ("누수는 없었나?",
         "초기 PoC의 보험 가입 수 변수(중요도 47.6%)가 라벨 정보를 재주입하고 있었다",
         "누수로 판정, 초기 AUROC 0.9998 자진 제외 — 새 검증은 외부 명단 포착률(5.2배)이라는 누수 불가능 지표로 교체"),
        ("통계법·데이터 변질 문제는?",
         "MDIS 승인 데이터의 가공·행정 활용은 통계 목적 외 사용 소지",
         "핵심 파이프라인을 MDIS 없이 전면 공개 데이터(국민연금·체불 명단·E-9 현황)로 재구축 완료 — 제약 원천 해소, MDIS는 보조 연구용으로만"),
    ]
    y = Inches(2.0)
    ch, gy = Inches(1.28), Inches(0.22)
    for i, (q, risk, ans) in enumerate(rows):
        yy = y + i * (ch + gy)
        rect(s, Inches(0.9), yy, Inches(11.6), ch, CARD, line=CARD_LINE,
             line_w=Pt(1), rounded=True)
        rect(s, Inches(0.9), yy, Inches(3.3), ch, NAVY, rounded=False)
        txt(s, Inches(1.15), yy, Inches(2.9), ch, [[(q, 15, WHITE, True)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        txt(s, Inches(4.45), yy + Inches(0.12), Inches(7.85), ch - Inches(0.2),
            [[("위험  ", 11.5, CORAL, True), (risk, 11.5, GRAY, False)],
             [("대응  ", 12, TEAL, True), (ans, 12, INK, True)]],
            space_after=Pt(5), line_spacing=1.08)
    yb = Inches(6.55)
    txt(s, Inches(0.9), yb, Inches(11.6), Inches(0.45),
        [[("약점을 먼저 공개하는 이유 — 정책 도구는 화려한 수치보다 검증 가능한 정직함이 먼저이기 때문입니다.",
           13, NAVY, True)]])


def s11_efficiency(prs):
    s = add_slide(prs)
    header(s, "05  RESULTS", "점검 효율 실측 — 고용 기준 5.2배 · 사업장 기준 2.4배", 11)
    im = img("fig_cell_capture.png")
    if im:
        rect(s, Inches(0.9), Inches(1.95), Inches(6.9), Inches(4.55), CARD,
             line=CARD_LINE, line_w=Pt(1), rounded=True)
        s.shapes.add_picture(im, Inches(1.05), Inches(2.15), width=Inches(6.6))
    # 우측 카드
    x = Inches(8.05)
    stat_card(s, x, Inches(1.95), Inches(4.45), Inches(1.35),
              "5.2배 · 2.4배", "고용 기준 · 사업장 수 기준 (상위 10% 셀)",
              "무작위 대비 · 체불 명단 709명 매칭 실측", accent=TEAL)
    stat_card(s, x, Inches(3.45), Inches(4.45), Inches(1.35),
              "30 → 39%", "상위 10% / 20% 셀 체불 포착률",
              "고용 6% / 10%만 점검하고도", accent=AMBER)
    rect(s, x, Inches(4.95), Inches(4.45), Inches(1.55), NAVY, rounded=True)
    txt(s, x + Inches(0.3), Inches(5.1), Inches(3.9), Inches(1.3),
        [[("시뮬레이션이 아니라 실측", 14, WHITE, True)],
         [("순열검정 p=0.006 · 건설업 제외에도 유지 —", 12.5, LGRAY, False)],
         [("시군구 해상도(1,873셀)에서도 4.9배(고용)", 12.5, TEAL, True)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(5), line_spacing=1.05)


def s12_riskmap(prs):
    s = add_slide(prs)
    header(s, "05  RESULTS", "시도 × 산업 위험 지도 — 어디부터 점검하나 (실측)", 12)
    im = img("fig_cell_riskmap.png")
    if im:
        rect(s, Inches(0.9), Inches(1.95), Inches(7.5), Inches(4.55), CARD,
             line=CARD_LINE, line_w=Pt(1), rounded=True)
        s.shapes.add_picture(im, Inches(1.02), Inches(2.35), width=Inches(7.26))
    x = Inches(8.65)
    txt(s, x, Inches(2.0), Inches(3.85), Inches(0.6),
        [[("고위험 상위 셀 (백분위)", 15, NAVY, True)]])
    clusters = [
        ("건설업 · 세종/제주/강원", "98~100", CORAL),
        ("숙박음식업 · 충북 등", "90~96", AMBER),
        ("농림어업 · 전남/경북", "88~94", AMBER),
    ]
    y = Inches(2.65)
    for i, (t, v, col) in enumerate(clusters):
        yy = y + i * Inches(0.78)
        rect(s, x, yy, Inches(3.85), Inches(0.65), CARD, line=CARD_LINE,
             line_w=Pt(1), rounded=True)
        rect(s, x, yy, Inches(0.12), Inches(0.65), col)
        txt(s, x + Inches(0.3), yy, Inches(2.5), Inches(0.65),
            [[(t, 12.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(2.7), yy, Inches(1.0), Inches(0.65),
            [[(v, 13.5, col, True)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    txt(s, x, Inches(5.15), Inches(3.85), Inches(1.5),
        [[("고위험 셀 = 실제 체불 명단 최다 셀과 일치", 12, INK, True)],
         [("(건설업 282명 · 제조업 186명 · 숙박음식 58명)", 11.5, GRAY, False)],
         [("농림어업·숙박음식은 E-9 집중 업종 — 이주노동자 보호 우선 지역", 11.5, GRAY, False)]],
        space_after=Pt(6), line_spacing=1.15)


def s12b_migrant(prs):
    s = add_slide(prs)
    header(s, "05  RESULTS", "이주노동자 보호 우선순위 — 노출-가중 실측", 13)
    im = img("fig_migrant_priority.png")
    if im:
        rect(s, Inches(0.9), Inches(1.95), Inches(7.5), Inches(4.55), CARD,
             line=CARD_LINE, line_w=Pt(1), rounded=True)
        s.shapes.add_picture(im, Inches(1.02), Inches(2.35), width=Inches(7.26))
    x = Inches(8.65)
    txt(s, x, Inches(2.0), Inches(3.85), Inches(0.9),
        [[("E-9 285,152명 (2026.1분기)", 15, NAVY, True)],
         [("KOSIS 시도×업종 실측 결합", 11.5, GRAY, False)]],
        space_after=Pt(3))
    rows = [
        ("규모의 축", "경기 제조업 — E-9 92,864명", NAVY),
        ("농도의 축", "전남·경기·충남 농림어업\n위험 백분위 84~97", CORAL),
        ("행정 액션", "통역 동반 점검·모국어 안내를\n이 순서로 배정", TEAL),
    ]
    y = Inches(3.0)
    for i, (t, d, col) in enumerate(rows):
        yy = y + i * Inches(1.02)
        rect(s, x, yy, Inches(3.85), Inches(0.9), CARD, line=CARD_LINE,
             line_w=Pt(1), rounded=True)
        rect(s, x, yy, Inches(0.12), Inches(0.9), col)
        txt(s, x + Inches(0.28), yy + Inches(0.1), Inches(3.45), Inches(0.75),
            [[(t, 12, col, True)]] +
            [[(ln, 11, INK, False)] for ln in d.split("\n")],
            space_after=Pt(2), line_spacing=1.02)
    txt(s, x, Inches(6.2), Inches(3.85), Inches(0.5),
        [[("우선순위 = E-9 인원 × 셀 위험 백분위 — 두 실측의 곱, 임의 가중 없음",
           10.5, GRAY, False)]], line_spacing=1.1)


def s13_impact(prs):
    s = add_slide(prs)
    header(s, "06  OPERATION", "누가 쓰고, 어떻게 돌아가는가 — 월간 운영 사이클", 14)
    # 사용자 3주체
    users = [
        ("근로감독관 (지청)", "매월 관할 우선순위 확인\n점검 계획서 발급 → 현장 점검\n결과 입력 → 명중률 환류", NAVY),
        ("고용노동부 본부·지방청", "전국 위험 지도로 감독 계획 수립\n분기 감독 자원 배분\n정책 효과 정량 모니터링", TEAL),
        ("이주노동자 지원기관", "E-9 밀집 × 고위험 셀 교집합에\n통역 동반 사전 안내·상담 배정", AMBER),
    ]
    x0, y = Inches(0.9), Inches(1.95)
    cw, ch, gap = Inches(3.75), Inches(1.9), Inches(0.22)
    for i, (t, d, col) in enumerate(users):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
        rect(s, x, y, cw, Inches(0.5), col)
        txt(s, x + Inches(0.2), y, cw - Inches(0.4), Inches(0.5),
            [[(t, 13.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.25), y + Inches(0.62), cw - Inches(0.45),
            ch - Inches(0.7),
            [[(ln, 11, INK, False)] for ln in d.split("\n")],
            space_after=Pt(3), line_spacing=1.05)
    # 월간 탐지 사이클 타임라인
    y2 = Inches(4.15)
    txt(s, Inches(0.9), y2, Inches(11.6), Inches(0.4),
        [[("탐지 사이클 — 신고를 기다리지 않는다", 15, NAVY, True)]])
    steps = [
        ("매월 말", "국민연금 사업장\n데이터 공표", NAVY),
        ("+1일", "자동 수집·재계산\n(스케줄러)", TEAL),
        ("+1일", "관할별 우선순위 갱신\n위험 급등 셀 경보", TEAL),
        ("월중", "감독관 표적 점검\n지원기관 사전 안내", AMBER),
        ("상시", "결과 입력 → 명중률\n→ 지표 재보정", CORAL),
    ]
    xs, ys = Inches(0.9), y2 + Inches(0.55)
    cw2, ch2, gap2 = Inches(2.14), Inches(1.35), Inches(0.22)
    for i, (tag, body, col) in enumerate(steps):
        x = xs + i * (cw2 + gap2)
        rect(s, x, ys, cw2, ch2, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
        rect(s, x, ys, cw2, Inches(0.38), col)
        txt(s, x + Inches(0.12), ys, cw2 - Inches(0.24), Inches(0.38),
            [[(tag, 11.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.14), ys + Inches(0.48), cw2 - Inches(0.28),
            Inches(0.85),
            [[(ln, 10.5, INK, False)] for ln in body.split("\n")],
            align=PP_ALIGN.CENTER, space_after=Pt(2), line_spacing=1.0)
        if i < len(steps) - 1:
            arr = s.shapes.add_shape(
                MSO_SHAPE.CHEVRON, x + cw2 + Emu(int(gap2) // 2) - Inches(0.13),
                ys + ch2 / 2 - Inches(0.13), Inches(0.26), Inches(0.26))
            _set_fill(arr, TEAL)
    txt(s, Inches(0.9), Inches(6.35), Inches(11.6), Inches(0.5),
        [[("주기 — 공개 데이터 기준 월 단위 자동 탐지. ", 12, INK, True),
          ("체불은 수개월에 걸쳐 악화되는 과정이라 월 주기로 선행 포착 가능 · "
           "기관 연계(고용보험 취득·상실 신고) 시 일 단위 준실시간으로 확장",
           11.5, GRAY, False)]], line_spacing=1.15)


def s14_roadmap(prs):
    s = add_slide(prs)
    header(s, "06  ROADMAP", "단계적 실증 로드맵", 15)
    txt(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(0.6),
        [[("이미 작동하는 파이프라인 위에서 ", 15, INK, False),
          ("수상 후 곧바로 실증 확장", 15, TEAL, True), (" 가능", 15, INK, False)]])
    phases = [
        ("1단계", "M+1~3", "외부 검증·고도화",
         "임금체불 명단·근로사건 매칭\n누수 제거 재학습·실측 가중치\n개방판 데이터 전환 착수", NAVY),
        ("2단계", "M+4~6", "광역 시범 적용",
         "1개 광역 지자체·지청 협력\n상위 50~100개 사전 점검\n점검 결과 모델 환류", TEAL),
        ("3단계", "M+7~12", "전국 확산·시스템 연계",
         "노동포털·EPS·고용24 연계\n지원기관 4~6곳 연동\n자동 재학습 스케줄러", AMBER),
    ]
    x0, y = Inches(0.9), Inches(2.75)
    cw, ch, gap = Inches(3.72), Inches(3.4), Inches(0.22)
    for i, (p, term, title, body, col) in enumerate(phases):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
        rect(s, x, y, cw, Inches(0.95), col, rounded=False)
        txt(s, x + Inches(0.25), y + Inches(0.12), cw - Inches(0.5), Inches(0.4),
            [[(p, 20, WHITE, True), ("   " + term, 13, RGBColor(0xCB,0xD5,0xE1), False)]])
        txt(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.5), Inches(0.4),
            [[(title, 14, WHITE, True)]])
        txt(s, x + Inches(0.3), y + Inches(1.2), cw - Inches(0.55), ch - Inches(1.3),
            [[("•  " + ln, 12.5, INK, False)] for ln in body.split("\n")],
            space_after=Pt(10), line_spacing=1.05)


def s14b_limits(prs):
    s = add_slide(prs)
    header(s, "06  LIMITS", "한계를 정직하게 — 그리고 다음 단계", 16)
    cards = [
        ("명단의 대표성", "체불 명단은 고액·상습 확정자 위주의 선별 표본이다",
         "근로감독 결과·진정 통계 등 추가 지표로 검증 확장 (로드맵 1단계)", NAVY),
        ("셀 단위의 해상도", "셀 안의 개별 사업장 위험까지는 구분하지 못한다",
         "감독관 행정데이터와 결합하는 시범 적용 단계에서 정밀화", TEAL),
        ("시계열 미검증", "LNN 조기 경보는 합성 데이터 시뮬레이션 결과다",
         "고용보험 이력 등 실제 시계열 확보 후 재검증 (로드맵 1단계)", AMBER),
        ("분류 크로스워크", "국민연금 업종코드(구분류)→현행 대분류 매핑은 근사치다",
         "사업자번호 기반 정밀 연계는 기관 협력 단계에서 보정", CORAL),
    ]
    x0, y0 = Inches(0.9), Inches(2.05)
    cw, ch, gx, gy = Inches(5.66), Inches(1.95), Inches(0.28), Inches(0.25)
    for i, (t, limit, plan, col) in enumerate(cards):
        x = x0 + (i % 2) * (cw + gx)
        y = y0 + (i // 2) * (ch + gy)
        rect(s, x, y, cw, ch, CARD, line=CARD_LINE, line_w=Pt(1), rounded=True)
        rect(s, x, y, Inches(0.12), ch, col)
        txt(s, x + Inches(0.32), y + Inches(0.18), cw - Inches(0.6), ch - Inches(0.3),
            [[(t, 14.5, col, True)],
             [(limit, 11.5, GRAY, False)],
             [("→ " + plan, 12, INK, True)]],
            space_after=Pt(5), line_spacing=1.08)
    txt(s, Inches(0.9), Inches(6.5), Inches(11.6), Inches(0.45),
        [[("한계를 아는 시스템이 현장에서 신뢰받는 시스템입니다.", 13, NAVY, True)]])


def s15_closing(prs):
    s = add_slide(prs)
    bg(s, NAVY)
    rect(s, Inches(0.9), Inches(1.2), Inches(1.5), Pt(4), TEAL)
    txt(s, Inches(0.9), Inches(1.5), Inches(11.6), Inches(1.4),
        [[("신고할 수 없는 피해자를,", 34, WHITE, True)],
         [("국가가 ", 34, WHITE, True), ("먼저 찾아가도록", 34, TEAL, True)]],
        line_spacing=1.15)
    # 3 요약 포인트
    pts = [
        ("선제 스크리닝", "신고 이전에 공공데이터로 고위험 셀 발굴 · 감독 우선순위 제안"),
        ("실측 검증", "독립 명단 5.2배(고용)·2.4배(사업장) · 홀드아웃 4.7배 · p=0.006"),
        ("법적 청정", "승인 불필요 공개 데이터만으로 재구축 · 즉시 실증 가능"),
    ]
    x0, y = Inches(0.9), Inches(3.4)
    cw, ch, gap = Inches(3.72), Inches(1.6), Inches(0.22)
    for i, (t, d) in enumerate(pts):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, NAVY2, rounded=True)
        rect(s, x, y, cw, Inches(0.1), TEAL)
        txt(s, x + Inches(0.25), y + Inches(0.25), cw - Inches(0.5), Inches(0.5),
            [[(t, 17, TEAL, True)]])
        txt(s, x + Inches(0.25), y + Inches(0.75), cw - Inches(0.5), Inches(0.8),
            [[(d, 12.5, RGBColor(0xE2,0xE8,0xF0), False)]], line_spacing=1.1)
    # 윤리 + 감사
    txt(s, Inches(0.9), Inches(5.5), Inches(11.6), Inches(0.6),
        [[("윤리 원칙 — ", 13, TEAL, True),
          ("개인 비식별 통계만 사용 · 법적 기준 라벨 · 자동 처벌 금지, 점검 우선순위 제안으로만 활용",
           12.5, RGBColor(0xCB,0xD5,0xE1), False)]])
    rect(s, Inches(0.9), Inches(6.25), Inches(11.6), Pt(1.5), NAVY2)
    txt(s, Inches(0.9), Inches(6.45), Inches(11.6), Inches(0.6),
        [[("감사합니다.  ", 18, WHITE, True),
          ("선문대학교 컴퓨터공학과 김태성  ·  github.com/kts6450/WageGuard",
           13, LGRAY, False)]])


def main():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    for fn in [s01_title, s02_agenda, s03_problem, s04_gap, s05_idea, s06_data,
               s08_ai, s09_product, s10_results, s10b_defense,
               s11_efficiency, s12_riskmap, s12b_migrant, s13_impact,
               s14_roadmap, s14b_limits, s15_closing]:
        fn(prs)
    out = os.path.join(HERE, "WageGuard_발표.pptx")
    try:
        prs.save(out)
    except PermissionError:
        out = os.path.join(HERE, "WageGuard_발표_실증반영.pptx")
        prs.save(out)
        print("(기존 파일이 열려 있어 새 이름으로 저장했습니다)")
    print("saved:", out, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
