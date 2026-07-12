# -*- coding: utf-8 -*-
"""PPTX 레이아웃 감사 — 폰트 메트릭 기반 텍스트 오버플로 추정.

PowerPoint 렌더러 없이, 맑은 고딕 실측 폭으로 각 텍스트 프레임의
예상 렌더 높이를 계산해 박스 높이 초과(넘침) 후보를 찾는다.
"""
from __future__ import annotations

import sys

from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REG = r"C:\Windows\Fonts\malgun.ttf"
BOLD = r"C:\Windows\Fonts\malgunbd.ttf"
_cache = {}


def font(size_pt: float, bold: bool):
    key = (round(size_pt), bold)
    if key not in _cache:
        _cache[key] = ImageFont.truetype(BOLD if bold else REG,
                                         int(round(size_pt * 96 / 72)))
    return _cache[key]


def text_w_px(t: str, size_pt: float, bold: bool) -> float:
    return font(size_pt, bold).getlength(t)


def main():
    prs = Presentation("WageGuard_발표.pptx")
    emu_per_px = 914400 / 96
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            box_w = sh.width / emu_per_px
            box_h = sh.height / emu_per_px
            if box_w <= 0 or box_h <= 0:
                continue
            total_h = 0.0
            preview = ""
            for p in tf.paragraphs:
                runs = [(r.text, (r.font.size.pt if r.font.size else 12),
                         bool(r.font.bold)) for r in p.runs if r.text]
                if not runs:
                    continue
                line_w = sum(text_w_px(t, s, b) for t, s, b in runs)
                max_pt = max(s for _, s, _ in runs)
                ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.05
                n_lines = max(1, -(-int(line_w) // max(int(box_w), 1)))
                total_h += n_lines * max_pt * (96 / 72) * ls * 1.18
                sa = p.space_after.pt if p.space_after else 4
                total_h += sa * 96 / 72
                if not preview:
                    preview = runs[0][0][:28]
            if total_h > box_h * 1.12 and total_h - box_h > 12:
                issues.append((si, preview, round(box_h), round(total_h)))
    if issues:
        print("오버플로 의심 (슬라이드 | 시작 텍스트 | 박스px | 추정px):")
        for it in sorted(issues):
            print(f"  s{it[0]:02d} | {it[1]} | {it[2]} | {it[3]}")
    else:
        print("오버플로 의심 없음")


if __name__ == "__main__":
    main()
